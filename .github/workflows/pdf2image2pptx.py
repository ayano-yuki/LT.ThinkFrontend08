# scripts/pdf2image2pptx_stream.py
import sys
import os
import io
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm

def pdf_to_pptx_stream(pdf_path, pptx_path, dpi=150, slide_ratio="16:9"):
    if not os.path.isfile(pdf_path):
        raise FileNotFoundError(f"PDFファイルが見つかりません: {pdf_path}")

    # プレゼンテーション作成
    prs = Presentation()
    if slide_ratio == "16:9":
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    elif slide_ratio == "3:4":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        raise ValueError("slide_ratioは '16:9' か '3:4' のみ指定可能です")

    # 総ページ数取得
    first_page_image = convert_from_path(pdf_path, dpi=dpi, first_page=1, last_page=1)[0]
    total_pages = len(convert_from_path(pdf_path, dpi=dpi, first_page=1, last_page=1))

    # ページごとに処理
    for page_num in tqdm(range(1, total_pages + 1), desc="PDF→PPTX変換"):
        images = convert_from_path(pdf_path, dpi=dpi, first_page=page_num, last_page=page_num)
        img = images[0]

        # 空白スライド追加
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # メモリ上に保存して貼り付け
        img_stream = io.BytesIO()
        img.save(img_stream, format="PNG")
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # メモリ解放
        img.close()
        img_stream.close()

    # PPTX保存
    prs.save(pptx_path)
    print(f"PPTXを保存しました: {pptx_path}")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("使い方: python pdf2image2pptx_stream.py input.pdf output.pptx")
        sys.exit(1)

    pdf_path = sys.argv[1]
    pptx_path = sys.argv[2]

    try:
        pdf_to_pptx_stream(pdf_path, pptx_path)
    except Exception as e:
        print(f"エラー: {e}")
        sys.exit(1)
